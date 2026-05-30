from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Design Tokens ─────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x0F, 0x25, 0x44)   # primary dark
ORANGE  = RGBColor(0xE8, 0x70, 0x1A)   # accent
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF7, 0xF8, 0xFA) # card / panel bg
GRAY_LT = RGBColor(0xE8, 0xEB, 0xF0)   # divider / border
GRAY_MD = RGBColor(0x8A, 0x96, 0xA8)   # secondary text
GRAY_DK = RGBColor(0x2C, 0x3E, 0x55)   # body text

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

# ── Primitives ────────────────────────────────────────────────────────────────

def bg(slide, color=WHITE):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=None, line=None, lw=0.75):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else:    s.fill.background()
    if line: s.line.color.rgb = line; s.line.width = Pt(lw)
    else:    s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h,
        size=16, bold=False, italic=False,
        color=GRAY_DK, align=PP_ALIGN.LEFT, wrap=True):
    if not text: return
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = "Calibri"
    return tb

def multiline_txt(slide, lines, l, t, w, h, size=13, color=GRAY_DK,
                  line_spacing_pt=6, bold_first=False):
    """lines = list of strings; first line optionally bold."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        from pptx.util import Pt as Pt2
        p.space_after = Pt2(line_spacing_pt)
        r = p.add_run(); r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Calibri"
        r.font.bold = (bold_first and i == 0)
    return tb

# ── Reusable layout pieces ────────────────────────────────────────────────────

def content_slide_header(slide, title, subtitle=None, timecode=None):
    """White slide with left orange accent bar, large title."""
    # Thin left accent bar
    rect(slide, 0, 0, Inches(0.07), SLIDE_H, fill=ORANGE)
    # Thin top rule under title area
    rect(slide, Inches(0.25), Inches(1.32), Inches(12.8), Pt(1.5), fill=GRAY_LT)
    # Title
    txt(slide, title,
        Inches(0.28), Inches(0.22), Inches(11.5), Inches(1.0),
        size=36, bold=True, color=NAVY)
    if subtitle:
        txt(slide, subtitle,
            Inches(0.28), Inches(0.92), Inches(11.0), Inches(0.38),
            size=13, color=GRAY_MD)
    if timecode:
        txt(slide, timecode,
            Inches(11.4), Inches(0.28), Inches(1.8), Inches(0.38),
            size=10, color=GRAY_MD, align=PP_ALIGN.RIGHT)

def narration_strip(slide, text):
    """Slim narration footer."""
    rect(slide, 0, Inches(6.88), SLIDE_W, Inches(0.62), fill=NAVY)
    txt(slide, "▶  " + text,
        Inches(0.28), Inches(6.91), Inches(12.75), Inches(0.56),
        size=10.5, italic=True, color=RGBColor(0xCC, 0xDD, 0xFF), wrap=True)

def placeholder_img(slide, label, l, t, w, h):
    rect(slide, l, t, w, h, fill=OFF_WHITE, line=GRAY_LT, lw=1)
    # dashed feel: inner lighter rect
    pad = Inches(0.12)
    rect(slide, l+pad, t+pad, w-2*pad, h-2*pad, fill=None, line=GRAY_LT, lw=0.5)
    txt(slide, label,
        l, t + h/2 - Inches(0.25), w, Inches(0.5),
        size=12, color=GRAY_MD, align=PP_ALIGN.CENTER, italic=True)

def bullet(slide, items, lx, start_y, col_w, size=14, gap=Inches(0.55), color=GRAY_DK):
    y = start_y
    for item in items:
        # small orange square bullet
        rect(slide, lx, y + Pt(size*0.45), Pt(size*0.4), Pt(size*0.4), fill=ORANGE)
        txt(slide, item, lx + Inches(0.22), y, col_w - Inches(0.25), gap,
            size=size, color=color)
        y += gap
    return y


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Introduction  (0:00 – 0:20)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, WHITE)

# Left panel — dark navy
rect(s, 0, 0, Inches(5.6), SLIDE_H, fill=NAVY)

# Logo placeholder inside left panel
placeholder_img(s, "Logo / Institution", Inches(0.35), Inches(0.35), Inches(2.2), Inches(1.5))

# Orange accent strip
rect(s, 0, Inches(5.45), Inches(5.6), Inches(0.08), fill=ORANGE)

# Presenter info at bottom of left panel
txt(s, "[Full Name]",
    Inches(0.35), Inches(5.6), Inches(5.0), Inches(0.48),
    size=16, bold=True, color=WHITE)
txt(s, "[Role / Program]",
    Inches(0.35), Inches(6.06), Inches(5.0), Inches(0.38),
    size=13, color=RGBColor(0xAA, 0xCC, 0xFF))
txt(s, "[Institution]   •   2026",
    Inches(0.35), Inches(6.42), Inches(5.0), Inches(0.38),
    size=12, color=GRAY_MD)

# Right panel — white, big title
txt(s, "CEPAT",
    Inches(6.1), Inches(1.5), Inches(6.8), Inches(2.0),
    size=82, bold=True, color=ORANGE)

txt(s, "Cepat Emergency Planning\nand Action Tool",
    Inches(6.1), Inches(3.45), Inches(6.8), Inches(1.1),
    size=22, bold=False, color=NAVY, wrap=True)

rect(s, Inches(6.1), Inches(4.7), Inches(2.0), Pt(3), fill=ORANGE)

txt(s, "Faster, Structured Earthquake Response",
    Inches(6.1), Inches(4.85), Inches(6.8), Inches(0.5),
    size=15, italic=True, color=GRAY_MD)

narration_strip(s,
    "Hello, my name is [Full Name]. I am [Role] at [Institution]. Today I will present CEPAT — "
    "the Cepat Emergency Planning and Action Tool — a multi-agent system designed to help BPBD respond "
    "to earthquakes faster, with clear workflows, reliable data, and human approval at every step.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement  (0:20 – 0:55)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, WHITE)
content_slide_header(s, "Problem Statement",
                     subtitle="Current disaster-response workflows create dangerous gaps",
                     timecode="0:20 – 0:55")

# 5 pain-point cards (left side)
problems = [
    ("Manual data gathering",     "Staff manually collect BMKG updates and search news under time pressure."),
    ("Slow alert drafting",       "Reports and alerts are written by hand, causing delays and inconsistency."),
    ("Fragmented coordination",   "Resource decisions are scattered across channels with no single source of truth."),
    ("Limited traceability",      "No audit trail means decisions are hard to review or justify after the fact."),
    ("Misinformation risk",       "Unverified reports spread fast, making it hard to prioritize credible data."),
]
py = Inches(1.5)
for title, body in problems:
    rect(s, Inches(0.28), py, Inches(6.0), Inches(0.72), fill=OFF_WHITE)
    rect(s, Inches(0.28), py, Pt(3), Inches(0.72), fill=ORANGE)
    txt(s, title, Inches(0.48), py + Inches(0.06), Inches(5.7), Inches(0.3),
        size=12, bold=True, color=NAVY)
    txt(s, body, Inches(0.48), py + Inches(0.36), Inches(5.7), Inches(0.32),
        size=11, color=GRAY_MD)
    py += Inches(0.85)

# Right: workflow diagram placeholder
placeholder_img(s,
    "Diagram: Current Manual Workflow\n\nBMKG data  →  news checks  →  manual drafting  →  scattered decisions",
    Inches(6.65), Inches(1.45), Inches(6.4), Inches(5.0))

narration_strip(s,
    "After a major earthquake, responders need rapid, consistent decisions. Current workflows are often manual: "
    "staff gather official updates, search news, draft alerts, and coordinate resources under time pressure. "
    "This creates delays, inconsistent messaging, and limited traceability. CEPAT addresses these issues by "
    "automating early analysis while keeping humans firmly in control.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Methodology  (0:55 – 2:00)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, WHITE)
content_slide_header(s, "Methodology",
                     subtitle="Python · Flask · SQLite · Five coordinated agents",
                     timecode="0:55 – 2:00")

# ── Agent pipeline ─────────────────────────────────────────────────────────────
agents = [
    ("01", "Monitoring",     "Polls BMKG feed;\nstores events in DB"),
    ("02", "Intelligence",   "Collects RSS news;\nclassifies credibility"),
    ("03", "Analysis",       "Generates sitrep;\nassigns risk level"),
    ("04", "Communication",  "Drafts alerts in\nmultiple formats"),
    ("05", "Coordination",   "Proposes resources\n& prioritized actions"),
]

BW = Inches(2.22)
BH = Inches(2.3)
GAP = Inches(0.16)
SX  = Inches(0.28)
SY  = Inches(1.5)

for i, (num, name, desc) in enumerate(agents):
    x = SX + i * (BW + GAP)
    # Card: alternating tones
    fill_c = NAVY if i % 2 == 0 else GRAY_DK
    rect(s, x, SY, BW, BH, fill=fill_c)

    # Number — small, top right
    txt(s, num, x + BW - Inches(0.75), SY + Inches(0.1),
        Inches(0.65), Inches(0.38),
        size=11, bold=True, color=ORANGE, align=PP_ALIGN.RIGHT)

    # Agent name
    txt(s, name + "\nAgent",
        x + Inches(0.18), SY + Inches(0.42),
        BW - Inches(0.25), Inches(0.85),
        size=16, bold=True, color=WHITE)

    # Description
    txt(s, desc,
        x + Inches(0.18), SY + Inches(1.28),
        BW - Inches(0.25), Inches(0.9),
        size=11, color=RGBColor(0xBB, 0xCC, 0xDD), wrap=True)

    # Arrow connector (between boxes)
    if i < 4:
        ax = x + BW + Inches(0.02)
        txt(s, "›",  # › single right angle quotation mark, clean arrow
            ax, SY + Inches(0.92), GAP + Inches(0.08), Inches(0.48),
            size=22, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# ── Output bar ────────────────────────────────────────────────────────────────
rect(s, Inches(0.28), Inches(3.98), Inches(12.77), Inches(0.04), fill=ORANGE)

outputs = [
    ("Flask Dashboard",  "Central output view"),
    ("Approval Queue",   "Human review gate"),
    ("Audit Log",        "Full decision trail"),
    ("SQLite Database",  "Local, no cloud needed"),
]
ox = Inches(0.28)
ow = Inches(3.1)
for t, sub in outputs:
    rect(s, ox, Inches(4.1), ow, Inches(0.82), fill=OFF_WHITE)
    txt(s, t,   ox+Inches(0.15), Inches(4.16), ow-Inches(0.2), Inches(0.35),
        size=12, bold=True, color=NAVY)
    txt(s, sub, ox+Inches(0.15), Inches(4.5),  ow-Inches(0.2), Inches(0.3),
        size=10.5, color=GRAY_MD)
    ox += ow + Inches(0.12)

narration_strip(s,
    "CEPAT is built in Python with a lightweight, modular architecture. Five coordinated agents reflect real "
    "response stages: (1) Monitoring polls BMKG; (2) Intelligence collects and classifies news; "
    "(3) Analysis generates a sitrep and risk level; (4) Communication drafts multi-format alerts; "
    "(5) Coordination proposes prioritized actions. All outputs appear in a Flask dashboard with an approval "
    "queue for human-in-the-loop review and an audit log for full accountability.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Results  (2:00 – 2:40)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, WHITE)
content_slide_header(s, "Results",
                     subtitle="Historical earthquake demo — complete pipeline in one automated flow",
                     timecode="2:00 – 2:40")

# Three screenshot placeholders (matching script visuals)
shots = ["Dashboard\n(main view)", "Situation Report\n& Alert Draft", "Approval Queue\n(human review)"]
px = Inches(0.28)
for label in shots:
    placeholder_img(s, label, px, Inches(1.48), Inches(4.1), Inches(3.3))
    px += Inches(4.3)

# Result bullets below screenshots
results = [
    "Automatically ingests official BMKG data",
    "Gathers and classifies related news reports by credibility",
    "Produces a structured situation report with risk level",
    "Drafts multi-format alerts ready for human approval",
    "Generates a prioritized coordination and resource plan",
]
bullet(s, results, Inches(0.35), Inches(4.98), Inches(12.5),
       size=13, gap=Inches(0.38))

narration_strip(s,
    "In a historical earthquake demo scenario, CEPAT completes a full response pipeline in one automated flow. "
    "It ingests official data, gathers related reports, produces a situation report, drafts alerts, "
    "and proposes a prioritized coordination plan — reducing manual steps and centralizing information.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Discussion  (2:40 – 3:20)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, WHITE)
content_slide_header(s, "Discussion",
                     subtitle="Impact and design rationale",
                     timecode="2:40 – 3:20")

# Left: responder screenshot
placeholder_img(s, "Screenshot / Photo:\nResponder using CEPAT dashboard",
                Inches(0.28), Inches(1.48), Inches(5.55), Inches(5.0))

# Right: three impact blocks
blocks = [
    ("Speed",
     "Automates data collection and first-pass analysis so responders focus on decisions, not repetitive tasks."),
    ("Transparency",
     "Every draft passes through a human approval queue. Every decision is recorded in an audit log."),
    ("Deployability",
     "Runs on standard local hardware — no cloud dependency. Suitable for regional agencies with limited resources."),
]
by = Inches(1.48)
for i, (title, body) in enumerate(blocks):
    rect(s, Inches(6.15), by, Inches(6.9), Inches(1.48), fill=OFF_WHITE)
    # orange top edge
    rect(s, Inches(6.15), by, Inches(6.9), Inches(0.055), fill=ORANGE)
    txt(s, title,
        Inches(6.35), by + Inches(0.12), Inches(6.5), Inches(0.42),
        size=15, bold=True, color=NAVY)
    txt(s, body,
        Inches(6.35), by + Inches(0.56), Inches(6.5), Inches(0.78),
        size=12, color=GRAY_DK, wrap=True)
    by += Inches(1.68)

narration_strip(s,
    "The main impact of CEPAT is speed with accountability. By automating data collection and first-pass analysis, "
    "responders can focus on decisions rather than repetitive tasks. Human approval and audit logging keep the system "
    "transparent and trustworthy. The architecture is lightweight and deployable on standard local infrastructure.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Conclusion  (3:20 – 3:50)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, WHITE)

# Left panel: dark, "thank you" side
rect(s, 0, 0, Inches(4.5), SLIDE_H, fill=NAVY)
rect(s, 0, 0, Inches(4.5), Inches(0.07), fill=ORANGE)

txt(s, "CEPAT",
    Inches(0.35), Inches(1.5), Inches(3.8), Inches(1.1),
    size=48, bold=True, color=ORANGE)

txt(s, "Cepat Emergency\nPlanning and\nAction Tool",
    Inches(0.35), Inches(2.6), Inches(3.8), Inches(1.5),
    size=16, color=WHITE, wrap=True)

rect(s, Inches(0.35), Inches(4.25), Inches(1.4), Pt(2.5), fill=ORANGE)

txt(s, "Thank you for watching.",
    Inches(0.35), Inches(4.42), Inches(3.8), Inches(0.5),
    size=14, italic=True, color=RGBColor(0xBB, 0xCC, 0xFF))

txt(s, "[Email or Website]",
    Inches(0.35), Inches(4.95), Inches(3.8), Inches(0.38),
    size=11, color=GRAY_MD)

txt(s, "3:20 – 3:50",
    Inches(0.35), Inches(6.95), Inches(3.5), Inches(0.35),
    size=9, color=GRAY_MD)

# Right: key takeaways
txt(s, "Key Takeaways",
    Inches(4.85), Inches(0.3), Inches(8.0), Inches(0.75),
    size=30, bold=True, color=NAVY)
rect(s, Inches(4.85), Inches(1.05), Inches(8.2), Pt(1.5), fill=GRAY_LT)

takeaways = [
    ("Faster Situational Awareness",
     "BMKG data and classified news signals processed automatically."),
    ("Structured Coordination",
     "A clear five-agent pipeline that mirrors real disaster response stages."),
    ("Reliable Communication",
     "Multi-format alert drafts ready for human review and release."),
    ("Human Oversight",
     "Approval queue and audit log ensure trust and accountability."),
    ("Anti-Misinformation",
     "Automatic credibility classification prioritizes verified information."),
]

ty = Inches(1.22)
for i, (title, body) in enumerate(takeaways):
    # zebra stripe
    fill_c = OFF_WHITE if i % 2 == 0 else WHITE
    rect(s, Inches(4.85), ty, Inches(8.2), Inches(0.95), fill=fill_c)
    # orange left rule
    rect(s, Inches(4.85), ty, Pt(3), Inches(0.95), fill=ORANGE)
    txt(s, title,
        Inches(5.08), ty + Inches(0.07), Inches(7.8), Inches(0.36),
        size=13, bold=True, color=NAVY)
    txt(s, body,
        Inches(5.08), ty + Inches(0.44), Inches(7.8), Inches(0.42),
        size=12, color=GRAY_MD)
    ty += Inches(1.05)

narration_strip(s,
    "To conclude, CEPAT delivers faster situational awareness, structured coordination, and reliable communication "
    "for earthquake response. It combines official data, news signals, and multi-agent analysis with strong human oversight. "
    "We believe CEPAT can strengthen local disaster response and reduce the impact of misinformation. Thank you.")


# ── Save ──────────────────────────────────────────────────────────────────────
out = r"C:\Users\MBN0C\Downloads\CEPAT\CEPAT_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
